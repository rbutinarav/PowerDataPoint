import streamlit as st
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


st.title('Power Data Point')

st.write("This is an example to show how a PPT can be updated with cool graphics easily.")

st.write("You would like to provide a ppt template ? Otherwise I will create a new one.")

#ask user to upload a ppt template
ppt_template_upload = st.file_uploader("Please upload a ppt template you would like to use.")

#ask user to define the slide to be used as a template, if nothing is selected, the last one will be used
slide_to_be_used_as_template = st.number_input("Please define the slide to be used as a template, if nothing is selected, the last one will be used", value=0)

#ask user ok to generate the ppt
ok_to_generate = st.button("Generate the ppt")

if ok_to_generate:
    st.write("I will generate a random dataframe and plot it.")

    df = pd.DataFrame(np.abs(np.random.randn(200, 3)), columns=['a', 'b', 'c']).round(2)

    # Determine the number of rows and columns for the table
    rows, cols = df.shape

    st.line_chart(df)

    st.write("This is a table of the same data.")

    st.dataframe(df)

    st.write("This is a bar chart of the same data.")

    st.bar_chart(df)

    if ppt_template_upload is not None:
        st.write("I will use the ppt template you provided.")
        prs = Presentation(ppt_template_upload)

    else:
        st.write("I will create a new ppt template.")
        # Create a presentation object
        prs = Presentation()

    # Define slide layout
    slide_layout = prs.slide_layouts[5]  # choose a slide layout (5 = title and content)

    # Add a slide for Line Chart
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Beautiful Line Chart"  # changed the title here

    # Define chart data for Line Chart
    chart_data = CategoryChartData()
    chart_data.categories = df.index
    chart_data.add_series('A', df['a'])  # plot 'a' series
    chart_data.add_series('B', df['b'])  # plot 'b' series
    chart_data.add_series('C', df['c'])  # plot 'c' series

    # Create Line chart on slide
    slide_width = 10  # typical slide width
    slide_height = 5.63  # typical slide height
    cx, cy = Inches(9), Inches(6)  # chart size
    x = Inches((slide_width - cx.inches) / 2)  # calculate x coordinate
    y = Inches((slide_height - cy.inches) / 2 + 1.18)  # calculate y coordinate

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    # Define maximum number of rows per slide (excluding the header)
    max_rows_per_slide = 15  # adjust this value according to your needs

    # Calculate the number of slides needed
    num_of_slides = -(-rows // max_rows_per_slide)  # equivalent to math.ceil(rows / max_rows_per_slide)

    for slide_num in range(num_of_slides):
        # Add a slide for each table
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Data Table - Part {}".format(slide_num + 1)

        # Calculate the number of rows for the current table
        rows_in_this_slide = min(max_rows_per_slide, rows - slide_num * max_rows_per_slide)

        # Add table to the slide
        slide_width = 10  # typical slide width
        slide_height = 5.63  # typical slide height
        width, height = Inches(9), Inches(4)  # table size
        x = Inches((slide_width - width.inches) / 2)  # calculate x coordinate
        y = Inches((slide_height - height.inches) / 2 + 0.7)  # calculate y coordinate, shifted 3 cm (approximately 1.18 inches) lower

        table = slide.shapes.add_table(rows_in_this_slide + 1, cols, x, y, width, height).table

        # Set column names in the table
        for i in range(cols):
            cell = table.cell(0, i)
            cell.text = df.columns[i]
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(13)
                    paragraph.alignment = PP_ALIGN.CENTER

        # Set values in the table
        for i in range(rows_in_this_slide):
            for j in range(cols):
                cell = table.cell(i + 1, j)
                cell.text = "{:.2f}".format(df.values[i + slide_num * max_rows_per_slide, j])
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        paragraph.alignment = PP_ALIGN.CENTER

        # Add a textbox at the bottom right of the slide
        left = Inches(9)   # Adjust these values as necessary
        top = Inches(6.5)  # Adjust these values as necessary
        width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        p.text = str(slide_num + 1)  # Set the text to the slide number

    # Save presentation
    prs.save('test.pptx')

    st.write("PowerPoint presentation has been created with the name 'test.pptx' in the current working directory.")
